# -*- coding: utf-8 -*-
"""
Created on Tue Sep 24 09:19:22 2019

@author: payam.bagheri
"""

import pandas as pd
import numpy as np
from os import path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
import pptx
from pptx.util import Inches
from pptx.dml.color import RGBColor

dir_path = path.dirname(path.dirname(path.abspath(__file__)))
print(dir_path)
#data = pd.read_excel(dir_path + '/0_input_data/2114_scorecard.xlsx')
prs = Presentation(dir_path + '/0_input_data/NQ Concept Testing 2019_FINAL v7_BREAKS.pptx')

breaks = pd.DataFrame(columns = range(24), index = range(30))

for j in range(24):
    brks = []
    for i in range(len(prs.slides[j].shapes)):
        #print('i is ', i, prs.slides[j].shapes[i].name)
        try:
            text = prs.slides[j].shapes[i].text_frame.text
            #print(text.split("||"))
            brk = text.split("||")
            
            brks.extend(brk)
            #print(brks)
        except AttributeError:
            pass
        
    print(len(brks))
    print(brks)
    breaks[j] = pd.Series(brks)
    
 
breaks.to_excel(dir_path + '/0_output/2114_gs_breaks.xlsx')

data = pd.read_csv(dir_path + '/0_input_data/2114_gs_lapsed.csv')
len(data['1'][data['1'] != '.'])
data.columns

heatmap_dat = pd.DataFrame(columns = range(5))

for i in range(24):
    nonils =  data[str(i)][data[str(i)] != '.']
    lennonils = len(data[str(i)][data[str(i)] != '.'])
    gooddat = pd.DataFrame(nonils.values.reshape(5,int(lennonils/5)).T)
    heatmap_dat = heatmap_dat.append(gooddat, ignore_index=True)
     
heatmap_dat.to_excel(dir_path + '/0_output/2114_heatmap_data_lapsed.xlsx')

skiprows = list(range(0,25))
heat = pd.read_excel(dir_path + '\\0_input_data\\Heatmap-2114.xls', skiprows=skiprows)
heat.head()
heat.columns
heat['Unnamed: 1'].head()

from xlrd import open_workbook
from xlutils.copy import copy

wb = open_workbook(dir_path + '\\0_input_data\\Heatmap-2114.xls', formatting_info=True)
sheet = wb.sheet_by_name("Total")
cell = sheet.cell(38, 1) # The first cell
print("cell.xf_index is", cell.xf_index)
fmt = wb.xf_list[cell.xf_index]
print("type(fmt) is", type(fmt))
print("Dumped Info:")
fmt.dump()


wb = open_workbook(dir_path + '\\0_input_data\\test.xls', formatting_info=True)
sheet = wb.sheet_by_name("Sheet1")
cell = sheet.cell(0,0) # The first cell
print("cell.xf_index is", cell.xf_index)
fmt = wb.xf_list[cell.xf_index]
print("type(fmt) is", type(fmt))
print("Dumped Info:")
fmt.dump()





from docx import Document
from docx.enum.text import WD_COLOR_INDEX

doc = Document(dir_path + '\\0_input_data\\Payam_bagheri.docx')

for i in range(1,13):
    phrase = str(sheet.cell(i,0))[6:-9]
    fmt = wb.xf_list[sheet.cell(i,0).xf_index]
    fmt.background.pattern_colour_index
    print(fmt.background.pattern_colour_index)
    for para in doc.paragraphs :
     start = para.text.find(phrase)
     if start > -1 :
      pre = para.text[:start]
      post = para.text[start+len(phrase):]
      para.text = pre
      para.add_run(phrase)
      para.runs[1].font.highlight_color = WD_COLOR_INDEX.YELLOW
      para.add_run(post)

doc.save(dir_path + '\\0_output\\t2.docx')







                


                    

